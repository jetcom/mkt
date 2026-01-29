from rest_framework.views import APIView
from rest_framework.response import Response
from rest_framework import status
from rest_framework.parsers import MultiPartParser, FormParser
from django.conf import settings
from questions.models import Question, Course
from questions.serializers import QuestionDetailSerializer
import json
import io
import difflib


class GenerateQuestionsView(APIView):
    """Generate questions using AI"""

    def post(self, request):
        provider = request.data.get('provider', 'claude')  # 'claude' or 'openai'
        content = request.data.get('content', '')  # Source material
        question_type = request.data.get('type', 'multipleChoice')
        count = min(int(request.data.get('count', 5)), 20)  # Max 20 at a time
        difficulty = request.data.get('difficulty', 'medium')
        examples = request.data.get('examples', [])  # Example questions for style
        course_id = request.data.get('course_id')  # For duplicate detection
        tag_name = request.data.get('tag_name')  # For duplicate detection

        if not content:
            return Response({'error': 'Content is required'}, status=status.HTTP_400_BAD_REQUEST)

        # Fetch existing questions for duplicate detection
        existing_questions = []
        if course_id:
            qs = Question.objects.filter(course_id=course_id, deleted_at__isnull=True)
            if tag_name:
                qs = qs.filter(tags__name=tag_name)
            existing_questions = list(qs.values_list('text', flat=True)[:100])  # Limit to 100 for prompt size

        try:
            print(f"[AI Generate] Provider: {provider}, Type: {question_type}, Count: {count}, Content length: {len(content)}, Existing: {len(existing_questions)}")

            if provider == 'claude':
                questions = self._generate_with_claude(content, question_type, count, difficulty, examples, existing_questions)
            elif provider == 'openai':
                questions = self._generate_with_openai(content, question_type, count, difficulty, examples, existing_questions)
            else:
                return Response({'error': 'Invalid provider'}, status=status.HTTP_400_BAD_REQUEST)

            # Post-generation duplicate filtering
            if existing_questions:
                original_count = len(questions)
                questions = self._filter_duplicates(questions, existing_questions)
                filtered_count = original_count - len(questions)
                if filtered_count > 0:
                    print(f"[AI Generate] Filtered {filtered_count} duplicate questions")

            print(f"[AI Generate] Success: {len(questions)} questions generated")
            return Response({'questions': questions})
        except Exception as e:
            import traceback
            print(f"[AI Generate] Error: {str(e)}")
            traceback.print_exc()
            return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

    def _filter_duplicates(self, new_questions, existing_texts, threshold=0.85):
        """Filter out questions that are too similar to existing ones."""
        unique_questions = []
        for q in new_questions:
            new_text = q.get('text', '').lower().strip()
            is_duplicate = False
            for existing_text in existing_texts:
                existing_lower = existing_text.lower().strip()
                # Use sequence matcher for similarity
                similarity = difflib.SequenceMatcher(None, new_text, existing_lower).ratio()
                if similarity >= threshold:
                    print(f"[Duplicate] Skipping question (similarity={similarity:.2f}): {new_text[:50]}...")
                    is_duplicate = True
                    break
            if not is_duplicate:
                unique_questions.append(q)
        return unique_questions

    def _build_prompt(self, content, question_type, count, difficulty, examples, existing_questions=None):
        type_instructions = {
            'multipleChoice': 'Create multiple choice questions with exactly 4 options (1 correct, 3 wrong).',
            'trueFalse': 'Create true/false questions.',
            'shortAnswer': 'Create short answer questions that can be answered in 1-2 sentences.',
            'longAnswer': 'Create long answer/essay questions that require detailed responses.',
        }

        example_text = ""
        if examples:
            example_text = "\n\nHere are example questions to match the style:\n"
            for i, ex in enumerate(examples[:3], 1):
                example_text += f"\nExample {i}:\n{ex}\n"

        # Add existing questions to avoid duplicates
        existing_text = ""
        if existing_questions:
            existing_text = "\n\nIMPORTANT: The following questions already exist. DO NOT create questions that are similar to these:\n"
            for i, eq in enumerate(existing_questions[:30], 1):  # Limit to 30 to save tokens
                truncated = eq[:200] + "..." if len(eq) > 200 else eq
                existing_text += f"- {truncated}\n"
            existing_text += "\nCreate DIFFERENT questions that cover other aspects of the material.\n"

        # Handle mixed question types
        if question_type == 'mixed':
            type_instruction = """Create a MIX of different question types. Include a variety of:
- Multiple choice questions (4 options: 1 correct, 3 wrong)
- True/false questions
- Short answer questions (1-2 sentence answers)

Aim for roughly: 40% multiple choice, 30% true/false, 30% short answer.
Each question MUST include a "question_type" field specifying its type."""

            json_format = """[
  {
    "text": "The question text",
    "question_type": "multipleChoice",
    "answer_data": {
      "correct": "The correct answer",
      "wrong": ["Wrong 1", "Wrong 2", "Wrong 3"]
    },
    "difficulty": "medium",
    "points": 2
  },
  {
    "text": "Is this statement true or false?",
    "question_type": "trueFalse",
    "answer_data": {
      "correct": true
    },
    "difficulty": "medium",
    "points": 1
  },
  {
    "text": "Explain briefly...",
    "question_type": "shortAnswer",
    "answer_data": {
      "solution": "The expected answer"
    },
    "difficulty": "medium",
    "points": 2
  }
]"""
        else:
            type_instruction = type_instructions.get(question_type, 'Create questions appropriate for the type specified.')
            json_format = f"""[
  {{
    "text": "The question text",
    "question_type": "{question_type}",
    "answer_data": {{
      // For multipleChoice:
      "correct": "The correct answer",
      "wrong": ["Wrong 1", "Wrong 2", "Wrong 3"]

      // For trueFalse:
      "correct": true or false

      // For shortAnswer/longAnswer:
      "solution": "The expected answer"
    }},
    "difficulty": "{difficulty}",
    "points": 2
  }}
]"""

        prompt = f"""Generate {count} {difficulty} difficulty questions based on the following content.

{type_instruction}
{existing_text}
Return your response as a JSON array with this structure:
{json_format}

Content to generate questions from:
---
{content}
---
{example_text}

Return ONLY the JSON array, no other text."""
        return prompt

    def _parse_json_response(self, response_text):
        """Extract and parse JSON from LLM response, handling markdown code blocks."""
        import re

        if not response_text or not response_text.strip():
            raise ValueError("Empty response from AI")

        text = response_text.strip()

        # Try to extract JSON from markdown code blocks
        code_block_match = re.search(r'```(?:json)?\s*([\s\S]*?)\s*```', text)
        if code_block_match:
            text = code_block_match.group(1).strip()

        # Try to find JSON array
        array_match = re.search(r'\[[\s\S]*\]', text)
        if array_match:
            text = array_match.group(0)

        try:
            return json.loads(text)
        except json.JSONDecodeError as e:
            # Include part of the response in error for debugging
            preview = response_text[:500] if len(response_text) > 500 else response_text
            raise ValueError(f"Failed to parse AI response as JSON: {e}. Response preview: {preview}")

    def _generate_with_claude(self, content, question_type, count, difficulty, examples, existing_questions=None):
        import anthropic

        if not settings.ANTHROPIC_API_KEY:
            raise ValueError("Anthropic API key not configured")

        client = anthropic.Anthropic(api_key=settings.ANTHROPIC_API_KEY)
        prompt = self._build_prompt(content, question_type, count, difficulty, examples, existing_questions)

        print(f"[Claude] Sending request, prompt length: {len(prompt)}")

        try:
            message = client.messages.create(
                model="claude-sonnet-4-20250514",
                max_tokens=4096,
                messages=[
                    {"role": "user", "content": prompt}
                ]
            )
        except Exception as api_error:
            print(f"[Claude] API error: {api_error}")
            raise ValueError(f"Claude API error: {str(api_error)}")

        print(f"[Claude] Response received, stop_reason: {message.stop_reason}")

        if not message.content:
            raise ValueError("Claude returned empty content")

        response_text = message.content[0].text
        print(f"[Claude] Response text length: {len(response_text) if response_text else 0}")

        if not response_text:
            raise ValueError("Claude returned empty text in response")

        return self._parse_json_response(response_text)

    def _generate_with_openai(self, content, question_type, count, difficulty, examples, existing_questions=None):
        from openai import OpenAI

        if not settings.OPENAI_API_KEY:
            raise ValueError("OpenAI API key not configured")

        client = OpenAI(api_key=settings.OPENAI_API_KEY)
        prompt = self._build_prompt(content, question_type, count, difficulty, examples, existing_questions)

        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "You are an expert educator who creates clear, well-structured exam questions. Always respond with valid JSON only."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.7,
        )

        response_text = response.choices[0].message.content
        return self._parse_json_response(response_text)


class ImproveQuestionView(APIView):
    """Improve an existing question using AI"""

    def post(self, request):
        provider = request.data.get('provider', 'claude')
        question = request.data.get('question', {})
        instruction = request.data.get('instruction', 'Improve this question for clarity and difficulty')

        if not question:
            return Response({'error': 'Question is required'}, status=status.HTTP_400_BAD_REQUEST)

        try:
            prompt = f"""Improve this exam question based on the following instruction: {instruction}

Current question:
Type: {question.get('question_type', 'unknown')}
Text: {question.get('text', '')}
Answer data: {json.dumps(question.get('answer_data', {}))}

Return your response as a JSON object with the same structure:
{{
  "text": "The improved question text",
  "answer_data": {{ ... }},
  "explanation": "Brief explanation of what was changed"
}}

Return ONLY the JSON object, no other text."""

            if provider == 'claude':
                result = self._call_claude(prompt)
            else:
                result = self._call_openai(prompt)

            return Response({'improved': result})
        except Exception as e:
            return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

    def _call_claude(self, prompt):
        import anthropic
        client = anthropic.Anthropic(api_key=settings.ANTHROPIC_API_KEY)
        message = client.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=2048,
            messages=[{"role": "user", "content": prompt}]
        )
        return json.loads(message.content[0].text)

    def _call_openai(self, prompt):
        from openai import OpenAI
        client = OpenAI(api_key=settings.OPENAI_API_KEY)
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "You are an expert educator."},
                {"role": "user", "content": prompt}
            ],
        )
        return json.loads(response.choices[0].message.content)


class ValidateQuestionView(APIView):
    """Validate a question for issues"""

    def post(self, request):
        question = request.data.get('question', {})

        if not question:
            return Response({'error': 'Question is required'}, status=status.HTTP_400_BAD_REQUEST)

        issues = []
        text = question.get('text', '')
        answer_data = question.get('answer_data', {})

        # Check for common issues
        if len(text) < 10:
            issues.append({'type': 'warning', 'message': 'Question text seems too short'})

        if question.get('question_type') == 'multipleChoice':
            wrong = answer_data.get('wrong', [])
            if len(wrong) < 3:
                issues.append({'type': 'error', 'message': 'Multiple choice needs at least 3 wrong answers'})
            if not answer_data.get('correct'):
                issues.append({'type': 'error', 'message': 'Missing correct answer'})

        if question.get('question_type') == 'trueFalse':
            if 'correct' not in answer_data:
                issues.append({'type': 'error', 'message': 'Missing true/false answer'})

        # Check for ambiguous language
        ambiguous_phrases = ['might be', 'could be', 'sometimes', 'usually', 'often']
        for phrase in ambiguous_phrases:
            if phrase in text.lower():
                issues.append({'type': 'warning', 'message': f'Potentially ambiguous phrase: "{phrase}"'})

        return Response({
            'valid': len([i for i in issues if i['type'] == 'error']) == 0,
            'issues': issues
        })


class GenerateVariantView(APIView):
    """Generate a variant of an existing question using AI"""

    def post(self, request):
        question_id = request.data.get('question_id')
        block_id = request.data.get('block_id')
        provider = request.data.get('provider', 'claude')
        target_type = request.data.get('target_type')  # None = same type, or specific type

        if not question_id:
            return Response({'error': 'question_id is required'}, status=status.HTTP_400_BAD_REQUEST)

        try:
            original = Question.objects.get(id=question_id)
        except Question.DoesNotExist:
            return Response({'error': 'Question not found'}, status=status.HTTP_404_NOT_FOUND)

        # Get other variants in the block for context
        existing_variants = []
        if original.block:
            existing_variants = list(
                original.block.questions.exclude(id=original.id)
                .values_list('text', flat=True)[:5]
            )

        # Determine output type
        output_type = target_type if target_type else original.question_type

        try:
            prompt = self._build_variant_prompt(original, existing_variants, target_type)

            if provider == 'claude':
                variant_data = self._call_claude(prompt)
            else:
                variant_data = self._call_openai(prompt)

            # Determine points based on type conversion
            points = original.points
            if target_type and target_type != original.question_type:
                # Adjust points for type conversion
                type_points = {
                    'multipleChoice': 1,
                    'trueFalse': 1,
                    'shortAnswer': 2,
                    'longAnswer': 4,
                }
                points = type_points.get(output_type, original.points)

            # Get next variant number for this block
            next_variant_num = (original.block.questions.count() + 1) if original.block else 1

            # Create the new question
            new_question = Question.objects.create(
                question_bank=original.question_bank,
                block=original.block,
                question_type=output_type,
                text=variant_data.get('text', ''),
                answer_data=variant_data.get('answer_data', {}),
                points=points,
                difficulty=original.difficulty,
                is_bonus=original.is_bonus,
                is_required=original.is_required,
                variant_number=next_variant_num,
            )

            # Copy tags from original
            new_question.tags.set(original.tags.all())

            return Response(QuestionDetailSerializer(new_question).data)

        except Exception as e:
            return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

    def _build_variant_prompt(self, original, existing_variants, target_type=None):
        type_instructions = {
            'multipleChoice': 'Create a multiple choice question with exactly 4 options (1 correct, 3 wrong). The answer_data should have "correct" (string) and "wrong" (array of 3 strings) keys.',
            'trueFalse': 'Create a true/false question. The answer_data should have "correct" as true or false (boolean).',
            'shortAnswer': 'Create a short answer question (1-2 sentence answer). The answer_data should have "solution" (string).',
            'longAnswer': 'Create a long answer/essay question (paragraph answer). The answer_data should have "solution" (string).',
        }

        existing_text = ""
        if existing_variants:
            existing_text = "\n\nOther existing variants (avoid duplicating these):\n"
            for i, text in enumerate(existing_variants, 1):
                existing_text += f"- {text[:200]}...\n" if len(text) > 200 else f"- {text}\n"

        # Determine if we're converting types
        output_type = target_type if target_type else original.question_type
        is_converting = target_type and target_type != original.question_type

        if is_converting:
            type_change_instruction = f"""
IMPORTANT: Convert this question to a {output_type} format.
{type_instructions.get(output_type, '')}

The new question should test the same concept but as a different question type."""
        else:
            type_change_instruction = f"""
Keep the same question type ({output_type}).
{type_instructions.get(output_type, '')}"""

        prompt = f"""Create a NEW variant of this exam question. The variant should test the same concept but with different wording, examples, or values.

Original question:
Type: {original.question_type}
Text: {original.text}
Answer data: {json.dumps(original.answer_data)}
{type_change_instruction}
{existing_text}

Requirements:
1. Test the same underlying concept/skill
2. Use different wording, numbers, examples, or scenarios
3. Maintain the same difficulty level

Return your response as a JSON object:
{{
  "text": "The new variant question text",
  "answer_data": {{ ... structure appropriate for {output_type} ... }}
}}

Return ONLY the JSON object, no other text."""
        return prompt

    def _call_claude(self, prompt):
        import anthropic

        if not settings.ANTHROPIC_API_KEY:
            raise ValueError("Anthropic API key not configured")

        client = anthropic.Anthropic(api_key=settings.ANTHROPIC_API_KEY)
        message = client.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=2048,
            messages=[{"role": "user", "content": prompt}]
        )
        return json.loads(message.content[0].text)

    def _call_openai(self, prompt):
        from openai import OpenAI

        if not settings.OPENAI_API_KEY:
            raise ValueError("OpenAI API key not configured")

        client = OpenAI(api_key=settings.OPENAI_API_KEY)
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "You are an expert educator creating exam question variants."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.8,
        )
        return json.loads(response.choices[0].message.content)


class ExtractFileContentView(APIView):
    """Extract text content from uploaded files (PPTX, PDF, TXT)"""
    parser_classes = [MultiPartParser, FormParser]

    def post(self, request):
        if 'file' not in request.FILES:
            return Response({'error': 'No file uploaded'}, status=status.HTTP_400_BAD_REQUEST)

        uploaded_file = request.FILES['file']
        filename = uploaded_file.name.lower()
        print(f"[ExtractFile] Received file: {uploaded_file.name}, size: {uploaded_file.size}")

        try:
            if filename.endswith('.pptx'):
                content = self._extract_pptx(uploaded_file)
            elif filename.endswith('.txt'):
                content = uploaded_file.read().decode('utf-8')
            elif filename.endswith('.md'):
                content = uploaded_file.read().decode('utf-8')
            else:
                return Response(
                    {'error': f'Unsupported file type. Supported: .pptx, .txt, .md'},
                    status=status.HTTP_400_BAD_REQUEST
                )

            print(f"[ExtractFile] Extracted {len(content)} chars from {uploaded_file.name}")
            return Response({
                'content': content,
                'filename': uploaded_file.name,
                'chars': len(content)
            })
        except Exception as e:
            import traceback
            print(f"[ExtractFile] Error: {str(e)}")
            traceback.print_exc()
            return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

    def _extract_pptx(self, file):
        from pptx import Presentation

        file_bytes = file.read()
        print(f"[PPTX] Read {len(file_bytes)} bytes from file")

        prs = Presentation(io.BytesIO(file_bytes))
        print(f"[PPTX] Presentation loaded, {len(prs.slides)} slides")

        slides_content = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []

            # Extract text from shapes
            for shape in slide.shapes:
                # Regular text shapes
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

                # Tables
                if hasattr(shape, "has_table") and shape.has_table:
                    for row in shape.table.rows:
                        row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_text:
                            slide_text.append(" | ".join(row_text))

            # Extract notes
            if hasattr(slide, "has_notes_slide") and slide.has_notes_slide:
                if slide.notes_slide and hasattr(slide.notes_slide, "notes_text_frame"):
                    notes_frame = slide.notes_slide.notes_text_frame
                    if notes_frame and hasattr(notes_frame, "text"):
                        notes = notes_frame.text.strip()
                        if notes:
                            slide_text.append(f"\n[Notes: {notes}]")

            if slide_text:
                slides_content.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_text))
                print(f"[PPTX] Slide {slide_num}: {len(slide_text)} text items")

        return "\n\n".join(slides_content)
